"""Constrained text extraction, chunking, embedding, and Firestore indexing."""

from __future__ import annotations

import hashlib
import io
import logging
import math
import multiprocessing
import re
import resource
import shlex
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from datetime import UTC, datetime
from types import SimpleNamespace
from typing import Any

from docx import Document as WordDocument
from google import genai
from google.genai import types
from pptx import Presentation
from pypdf import PdfReader

from app.config import Settings
from app.domain.enums import DocumentStatus
from app.domain.models import DocumentChunk, RetrievalCitation
from app.domain.repository import Repository
from app.retrieval.object_store import DocumentObjectStore

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class ExtractedUnit:
    text: str
    page: int | None = None
    slide: int | None = None


def sanitize_text(value: str, max_length: int | None = None) -> str:
    cleaned = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", " ", value)
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
    return cleaned[:max_length] if max_length else cleaned


def extract_units(data: bytes, extension: str, settings: Settings) -> list[ExtractedUnit]:
    if extension == ".pdf":
        reader = PdfReader(io.BytesIO(data))
        if reader.is_encrypted:
            raise ValueError("encrypted_pdf")
        if len(reader.pages) > settings.document_max_pages:
            raise ValueError("page_limit_exceeded")
        units = [
            ExtractedUnit(sanitize_text(page.extract_text() or ""), page=index + 1)
            for index, page in enumerate(reader.pages)
        ]
        meaningful = [unit for unit in units if len(unit.text) >= 20]
        if not meaningful:
            raise ValueError("image_only_pdf")
        return meaningful
    if extension == ".docx":
        document = WordDocument(io.BytesIO(data))
        blocks = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            blocks.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
        text = sanitize_text("\n".join(blocks))
        return [ExtractedUnit(text)] if text else []
    if extension == ".pptx":
        presentation = Presentation(io.BytesIO(data))
        if len(presentation.slides) > settings.document_max_pages:
            raise ValueError("slide_limit_exceeded")
        units: list[ExtractedUnit] = []
        for index, slide in enumerate(presentation.slides):
            text = sanitize_text(
                "\n".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text
                )
            )
            if text:
                units.append(ExtractedUnit(text=text, slide=index + 1))
        return units
    text = sanitize_text(data.decode("utf-8"))
    return [ExtractedUnit(text)] if text else []


def _extract_child(
    connection: Any,
    data: bytes,
    extension: str,
    max_pages: int,
) -> None:
    """Parse one untrusted document inside a bounded child process."""
    try:
        # Cloud Run uses Linux. macOS reserves a very large virtual address
        # range for the interpreter, so RLIMIT_AS cannot be lowered reliably in
        # local tests; the parent timeout still isolates those test parses.
        if sys.platform == "linux":
            memory_limit = 2 * 1024 * 1024 * 1024
            _, memory_hard = resource.getrlimit(resource.RLIMIT_AS)
            bounded_memory = (
                memory_limit
                if memory_hard == resource.RLIM_INFINITY
                else min(memory_limit, memory_hard)
            )
            resource.setrlimit(resource.RLIMIT_AS, (bounded_memory, memory_hard))
            _, cpu_hard = resource.getrlimit(resource.RLIMIT_CPU)
            bounded_cpu = 45 if cpu_hard == resource.RLIM_INFINITY else min(45, cpu_hard)
            resource.setrlimit(resource.RLIMIT_CPU, (bounded_cpu, cpu_hard))
        units = extract_units(
            data,
            extension,
            SimpleNamespace(document_max_pages=max_pages),  # type: ignore[arg-type]
        )
        connection.send(("ok", units))
    except BaseException as exc:  # the parent converts this to a sanitized failure
        connection.send(("error", type(exc).__name__, sanitize_text(str(exc), 120)))
    finally:
        connection.close()


def extract_units_constrained(
    data: bytes,
    extension: str,
    settings: Settings,
) -> list[ExtractedUnit]:
    context = multiprocessing.get_context("spawn")
    receive, send = context.Pipe(duplex=False)
    process = context.Process(
        target=_extract_child,
        args=(send, data, extension, settings.document_max_pages),
        daemon=True,
    )
    process.start()
    send.close()
    try:
        if not receive.poll(60):
            process.terminate()
            process.join(timeout=5)
            if process.is_alive():
                process.kill()
                process.join(timeout=5)
            raise TimeoutError("document_extraction_timeout")
        result = receive.recv()
    finally:
        receive.close()
    process.join(timeout=5)
    if process.is_alive():
        process.kill()
        process.join(timeout=5)
    if not result or result[0] != "ok":
        error_type = str(result[1]) if len(result) > 1 else "ExtractionError"
        error_message = str(result[2]) if len(result) > 2 else "document_extraction_failed"
        if error_type == "ValueError":
            raise ValueError(error_message)
        raise RuntimeError("document_extraction_failed")
    return list(result[1])


def chunk_units(
    units: list[ExtractedUnit], target_tokens: int, overlap_tokens: int
) -> list[ExtractedUnit]:
    # A whitespace token approximation keeps extraction deterministic and avoids
    # coupling persisted chunks to a model-specific tokenizer release.
    words: list[tuple[str, int | None, int | None]] = []
    for unit in units:
        words.extend((word, unit.page, unit.slide) for word in unit.text.split())
    if not words:
        return []
    stride = max(1, target_tokens - overlap_tokens)
    chunks: list[ExtractedUnit] = []
    for start in range(0, len(words), stride):
        window = words[start : start + target_tokens]
        if not window:
            break
        pages = [item[1] for item in window if item[1] is not None]
        slides = [item[2] for item in window if item[2] is not None]
        chunks.append(
            ExtractedUnit(
                text=" ".join(item[0] for item in window),
                page=min(pages) if pages else None,
                slide=min(slides) if slides else None,
            )
        )
        if start + target_tokens >= len(words):
            break
    return chunks


class EmbeddingService:
    def __init__(self, settings: Settings) -> None:
        self.settings = settings
        self._client: genai.Client | None = None

    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        return self._embed(texts, "RETRIEVAL_DOCUMENT")

    def embed_query(self, text: str) -> list[float]:
        return self._embed([text], "RETRIEVAL_QUERY")[0]

    def _embed(self, texts: list[str], task_type: str) -> list[list[float]]:
        if self.settings.env in {"local", "test"}:
            return [self._deterministic_embedding(text) for text in texts]
        client = self._client or genai.Client(
            vertexai=True,
            project=self.settings.project_id,
            location=self.settings.region,
        )
        self._client = client
        vectors: list[list[float]] = []
        for offset in range(0, len(texts), 20):
            response = client.models.embed_content(
                model=self.settings.embedding_model,
                contents=texts[offset : offset + 20],
                config=types.EmbedContentConfig(
                    task_type=task_type,
                    output_dimensionality=self.settings.embedding_dimensions,
                ),
            )
            vectors.extend([list(item.values) for item in response.embeddings or []])
        if len(vectors) != len(texts):
            raise RuntimeError("Embedding service returned an unexpected result count")
        return vectors

    def _deterministic_embedding(self, text: str) -> list[float]:
        values = [0.0] * self.settings.embedding_dimensions
        for token in re.findall(r"[\w'-]+", text.casefold()):
            digest = hashlib.sha256(token.encode()).digest()
            index = int.from_bytes(digest[:4], "big") % len(values)
            values[index] += 1.0 if digest[4] & 1 else -1.0
        norm = math.sqrt(sum(value * value for value in values)) or 1.0
        return [value / norm for value in values]


class DocumentIndexer:
    def __init__(
        self,
        repository: Repository,
        object_store: DocumentObjectStore,
        embeddings: EmbeddingService,
        settings: Settings,
    ) -> None:
        self.repository = repository
        self.object_store = object_store
        self.embeddings = embeddings
        self.settings = settings

    def process(self, room_id: str, version_id: str) -> int:
        version = self.repository.get_document_version(room_id, version_id)
        if version.status == DocumentStatus.READY:
            return version.chunk_count or 0
        version.status = DocumentStatus.INDEXING
        version.updated_at = datetime.now(UTC)
        self.repository.save_document_version(version)
        try:
            data = self.object_store.download(version.object_name)
            self._scan(data, version.extension)
            units = extract_units_constrained(data, version.extension, self.settings)
            character_count = sum(len(unit.text) for unit in units)
            if not units or character_count == 0:
                raise ValueError("no_extractable_text")
            if character_count > self.settings.document_max_characters:
                raise ValueError("character_limit_exceeded")
            extracted_chunks = chunk_units(
                units,
                self.settings.document_chunk_tokens,
                self.settings.document_chunk_overlap_tokens,
            )
            vectors = self.embeddings.embed_documents([chunk.text for chunk in extracted_chunks])
            chunks: list[DocumentChunk] = []
            for index, (chunk, vector) in enumerate(zip(extracted_chunks, vectors, strict=True)):
                chunk_id = hashlib.sha256(
                    f"{version.id}\0{index}\0{chunk.text}".encode()
                ).hexdigest()
                chunks.append(
                    DocumentChunk(
                        id=chunk_id,
                        room_id=room_id,
                        document_id=version.document_id,
                        version_id=version.id,
                        title=version.title,
                        version=version.version,
                        text=chunk.text,
                        page_start=chunk.page,
                        page_end=chunk.page,
                        slide_start=chunk.slide,
                        slide_end=chunk.slide,
                        embedding=vector,
                        expires_at=version.expires_at,
                    )
                )
            self.repository.delete_document_chunks(room_id, version.id)
            self.repository.save_document_chunks(chunks)
            version.status = DocumentStatus.READY
            version.error_code = None
            version.error_message = None
            version.character_count = character_count
            version.chunk_count = len(chunks)
            version.page_count = max((unit.page or 0 for unit in units), default=0) or None
            version.slide_count = max((unit.slide or 0 for unit in units), default=0) or None
            version.ready_at = datetime.now(UTC)
            version.updated_at = version.ready_at
            self.repository.save_document_version(version)
            document = self.repository.get_document(room_id, version.document_id)
            if document.pending_version_id == version.id:
                document.active_version_id = version.id
                document.pending_version_id = None
                document.updated_at = version.ready_at
                self.repository.save_document(document)
            logger.info("event=document_index_succeeded chunks=%d", len(chunks))
            return len(chunks)
        except Exception as exc:
            version.status = DocumentStatus.FAILED
            version.error_code = self._error_code(exc)
            version.error_message = "Document could not be indexed"
            version.updated_at = datetime.now(UTC)
            self.repository.save_document_version(version)
            logger.error(
                "event=document_index_failed error_code=%s error_type=%s",
                version.error_code,
                type(exc).__name__,
            )
            raise

    def _scan(self, data: bytes, extension: str) -> None:
        # The standard EICAR marker is rejected even in hermetic tests. Deployed
        # environments additionally fail closed through the configured ClamAV command.
        if b"EICAR-STANDARD-ANTIVIRUS-TEST-FILE" in data:
            raise ValueError("malware_detected")
        command = shlex.split(self.settings.document_malware_scan_command)
        if not self.settings.document_malware_scan_required:
            return
        with tempfile.NamedTemporaryFile(suffix=extension) as handle:
            handle.write(data)
            handle.flush()
            try:
                result = subprocess.run(
                    [*command, handle.name],
                    check=False,
                    capture_output=True,
                    timeout=120,
                )
            except (FileNotFoundError, subprocess.TimeoutExpired) as exc:
                raise RuntimeError("malware_scan_unavailable") from exc
            if result.returncode != 0:
                raise ValueError("malware_detected")

    @staticmethod
    def _error_code(error: Exception) -> str:
        value = sanitize_text(str(error), 80)
        return value if re.fullmatch(r"[a-z0-9_ -]+", value.casefold()) else "index_failed"


class DocumentRetrievalService:
    def __init__(
        self, repository: Repository, embeddings: EmbeddingService, settings: Settings
    ) -> None:
        self.repository = repository
        self.embeddings = embeddings
        self.settings = settings

    def search(self, room_id: str, version_ids: list[str], query: str) -> list[RetrievalCitation]:
        cleaned_query = sanitize_text(query, 500)
        if not cleaned_query or not version_ids:
            return []
        query_vector = self.embeddings.embed_query(cleaned_query)
        matches = self.repository.search_document_chunks(
            room_id,
            version_ids,
            query_vector,
            self.settings.document_retrieval_limit,
        )
        citations: list[RetrievalCitation] = []
        for chunk, distance in matches:
            if distance > self.settings.document_retrieval_max_distance:
                continue
            citations.append(
                RetrievalCitation(
                    document_id=chunk.document_id,
                    version_id=chunk.version_id,
                    title=chunk.title,
                    version=chunk.version,
                    excerpt=sanitize_text(chunk.text, 1400),
                    page_start=chunk.page_start,
                    page_end=chunk.page_end,
                    slide_start=chunk.slide_start,
                    slide_end=chunk.slide_end,
                    distance=distance,
                )
            )
        return citations
